"""Generate docs/araroopat_roundtrip.pptx — a schematic of the يذهب round trip.

Companion to docs/araroopat_end_to_end.md (Part 10). Every value on the slides
is a real measurement from outputs/tokenizers/araroopat_balanced and the
Phase-3 checkpoint; see the markdown doc for provenance.

    .venv/bin/python docs/araroopat_roundtrip_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from deck_kit import *  # noqa: F403 — palette + box/arrow/text primitives + audit
from deck_kit import (MSO_ANCHOR, W, H, arrow, audit, blank, box, footnote, text, title)

OUT = Path(__file__).with_name("araroopat_roundtrip.pptx")

# ==========================================================================
# Slide 1 — title
# ==========================================================================
prs = Presentation()
prs.slide_width, prs.slide_height = W, H

s = blank(prs)
box(s, 0, 0, 13.333, 7.5, fill=RGBColor(0x0F, 0x17, 0x2A), line=None,
    shape=MSO_SHAPE.RECTANGLE)
box(s, 0, 0, 13.333, 0.10, fill=TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)

t = s.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.9))
text(t, "AraRooPat — the full round trip", size=40, bold=True,
     color=WHITE, align=PP_ALIGN.LEFT)
t = s.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(11.5), Inches(0.8))
text(t, [("يذهب", {"size": 46, "bold": True, "color": RGBColor(0x2D, 0xD4, 0xBF)})],
     align=PP_ALIGN.LEFT)
t = s.shapes.add_textbox(Inches(0.95), Inches(3.35), Inches(11.0), Inches(0.6))
text(t, "one Arabic word → 2 morphological tokens → LLaMA-3.2-1B → back to Arabic",
     size=16, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)

teaser = [("يذهب", SLATE_L, 1.30), ("CAMeL analyze", AMBER_L, 1.90),
          ("[ROOT_ذهب] [PAT_يَ1ْ2َ3]", INDIGO_L, 2.90), ("[1, 354, 3298, 2]", INDIGO_L, 2.10),
          ("LLaMA-3.2-1B", PURPLE_L, 1.70), ("يذهب", GREEN_L, 1.30)]
x = 0.9
for i, (lbl, col, w) in enumerate(teaser):
    box(s, x, 4.45, w, 0.62, lbl, fill=col, line=None, radius=0.25,
        size=10.5, bold=True, color=RGBColor(0x1E, 0x29, 0x3B))
    x += w
    if i < len(teaser) - 1:
        arrow(s, x + 0.03, 4.66, 0.18, 0.20, "right", RGBColor(0x47, 0x55, 0x69))
        x += 0.24

t = s.shapes.add_textbox(Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.9))
text(t, ["Measured end-to-end on: tokenizer outputs/tokenizers/araroopat_balanced (vocab 3,784)  ·  "
         "model outputs/experiments/araroopat_3phase_with_sft/training/sft  ·  live CAMeL bridge",
         "Companion to docs/araroopat_end_to_end.md"],
     size=11, color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.LEFT, space=4)

# ==========================================================================
# Slide 2 — the master schematic (Part 10 of the markdown)
# ==========================================================================
s = blank(prs)
title(s, "The full round trip", "input → tokenizer → ids → model → decode → output. Shapes carry real values.")

BH = 0.82   # band box height

# ---- Band 1: tokenizer front-end -------------------------------------------
box(s, 0.55, 1.28, 0.30, BH + 0.55, "1", fill=TEAL, line=None,
    shape=MSO_SHAPE.RECTANGLE, size=15, bold=True, color=WHITE)
lbl = s.shapes.add_textbox(Inches(0.95), Inches(1.30), Inches(3.0), Inches(0.26))
text(lbl, "TOKENIZER  ·  pure Python + CAMeL", size=10, bold=True, color=TEAL,
     align=PP_ALIGN.LEFT)

row1 = [
    (1.00, 1.55, [("يذهب", {"size": 20, "bold": True, "color": WHITE})], SLATE, None),
    (2.75, 1.75, ["NFKC normalize", ("+ whitespace split", {"size": 9.5, "color": MUTED})], WHITE, TEAL),
    (4.80, 2.05, ["character-class runs", ("[alpha: يذهب]  1 run", {"size": 9.5, "color": MUTED})], WHITE, TEAL),
    (7.15, 2.30, ["CAMeL analyze  (subprocess)", ("root 'ذ.ه.ب'  ·  pattern 'يَ1ْ2َ3'", {"size": 9.5, "color": AMBER})], AMBER_L, AMBER),
    (9.75, 3.00, ["post-process → Analysis", ("root='ذهب'  pattern='يَ1ْ2َ3'  clitics=()", {"size": 9.5, "color": MUTED})], WHITE, TEAL),
]
for x, w, lines, fill, line in row1:
    box(s, x, 1.62, w, BH, lines, fill=fill, line=line or SLATE_L, size=12, bold=True)
for x in (2.60, 4.65, 7.00, 9.60):
    arrow(s, x, 1.92, 0.16, 0.22, "right", TEAL)

# ---- Band 2: vocabulary + ids ---------------------------------------------
box(s, 0.55, 2.92, 0.30, BH + 0.30, "2", fill=INDIGO, line=None,
    shape=MSO_SHAPE.RECTANGLE, size=15, bold=True, color=WHITE)
lbl = s.shapes.add_textbox(Inches(0.95), Inches(2.94), Inches(4.0), Inches(0.26))
text(lbl, "VOCABULARY LOOKUP  ·  3,784 entries", size=10, bold=True, color=INDIGO,
     align=PP_ALIGN.LEFT)

box(s, 1.00, 3.24, 3.05, BH, ["[ROOT_ذهب]  →  id 354",
                              ("root token, ids 138–3283", {"size": 9, "color": MUTED})],
    fill=INDIGO_L, line=INDIGO, size=12.5, bold=True, color=INDIGO)
box(s, 4.35, 3.24, 3.30, BH, ["[PAT_يَ1ْ2َ3]  →  id 3298",
                              ("pattern token, ids 3284–3783", {"size": 9, "color": MUTED})],
    fill=INDIGO_L, line=INDIGO, size=12.5, bold=True, color=INDIGO)
box(s, 7.95, 3.24, 4.80, BH,
    [("input_ids = [ 1, 354, 3298, 2 ]", {"size": 14, "bold": True, "font": MONO, "color": WHITE}),
     ("<s>  ·  ROOT  ·  PAT  ·  </s>            attention_mask = [1,1,1,1]", {"size": 9, "color": RGBColor(0xC7, 0xD2, 0xFE), "font": MONO})],
    fill=INDIGO, line=None)
arrow(s, 4.10, 3.54, 0.16, 0.22, "right", INDIGO)
arrow(s, 7.72, 3.54, 0.16, 0.22, "right", INDIGO)
arrow(s, 11.05, 2.50, 0.24, 0.60, "down", INDIGO)

# ---- Band 3: model ---------------------------------------------------------
box(s, 0.55, 4.28, 0.30, BH + 0.30, "3", fill=PURPLE, line=None,
    shape=MSO_SHAPE.RECTANGLE, size=15, bold=True, color=WHITE)
lbl = s.shapes.add_textbox(Inches(0.95), Inches(4.30), Inches(5.0), Inches(0.26))
text(lbl, "MODEL  ·  LLaMA-3.2-1B, unmodified apart from vocab size", size=10,
     bold=True, color=PURPLE, align=PP_ALIGN.LEFT)

row3 = [
    (1.00, 2.30, ["embed_tokens", ("(3784 × 2048) gather → [B,4,2048]", {"size": 9, "color": MUTED})]),
    (3.60, 2.55, ["16 decoder layers", ("RoPE · GQA 32/8 · SwiGLU", {"size": 9, "color": MUTED})]),
    (6.45, 2.45, ["lm_head  (tied)", ("2048 → 3,784 logits", {"size": 9, "color": MUTED})]),
    (9.20, 3.55, ["cross-entropy / log-softmax", ("training loss  ·  eval log-likelihood", {"size": 9, "color": MUTED})]),
]
for x, w, lines in row3:
    box(s, x, 4.60, w, BH, lines, fill=PURPLE_L, line=PURPLE, size=12, bold=True, color=PURPLE)
for x in (3.40, 6.25, 9.00):
    arrow(s, x, 4.90, 0.16, 0.22, "right", PURPLE)
arrow(s, 11.05, 4.12, 0.24, 0.42, "down", PURPLE)

# ---- Band 4: decode --------------------------------------------------------
box(s, 0.55, 5.72, 0.30, BH + 0.28, "4", fill=GREEN, line=None,
    shape=MSO_SHAPE.RECTANGLE, size=15, bold=True, color=WHITE)
lbl = s.shapes.add_textbox(Inches(0.95), Inches(5.74), Inches(5.0), Inches(0.26))
text(lbl, "DECODE  ·  state machine + 3-tier reconstruction", size=10, bold=True,
     color=GREEN, align=PP_ALIGN.LEFT)

box(s, 1.00, 6.02, 2.75, 0.78, ["354 → pending_root = 'ذهب'"], fill=GREEN_L, line=GREEN,
    size=11.5, bold=True, color=GREEN)
box(s, 4.05, 6.02, 4.55, 0.78, ["3298 → reconstruction[(354, 3298)] = 'يذهب'",
                                ("tier 1: O(1) lookup — 60,248 entries", {"size": 9, "color": MUTED})],
    fill=GREEN_L, line=GREEN, size=11.5, bold=True, color=GREEN)
box(s, 8.90, 6.02, 2.10, 0.78, [("يذهب", {"size": 20, "bold": True, "color": WHITE})],
    fill=SLATE, line=None)
box(s, 11.15, 6.02, 1.60, 0.78, [("exact", {"size": 13, "bold": True, "color": GREEN}),
                                 ("round trip ✓", {"size": 9, "color": MUTED})],
    fill=WHITE, line=GREEN)
arrow(s, 3.82, 6.30, 0.16, 0.22, "right", GREEN)
arrow(s, 8.68, 6.30, 0.16, 0.22, "right", GREEN)
arrow(s, 11.05, 5.50, 0.24, 0.45, "down", GREEN)

footnote(s, "4 tokens for a 4-character word  ·  the model never sees the surface string, only the (root, pattern) pair")

# ==========================================================================
# Slide 3 — inside the tokenizer
# ==========================================================================
s = blank(prs)
title(s, "Inside the tokenizer — analysis crosses a process boundary",
      "camel-tools pins numpy<2 / transformers<4.54, so it runs in .venv-camel and is reached over NDJSON pipes.")

box(s, 0.55, 1.45, 2.25, 1.05, [("يذهب", {"size": 22, "bold": True, "color": WHITE}),
                                ("whitespace word", {"size": 9, "color": RGBColor(0xCB, 0xD5, 0xE1)})],
    fill=SLATE, line=None)
arrow(s, 2.92, 1.86, 0.22, 0.24, "right", TEAL)

box(s, 3.28, 1.30, 4.35, 1.35,
    [("main .venv", {"size": 10, "bold": True, "color": TEAL}),
     ("MorphAnalyzer.analyze(\"يذهب\")", {"size": 11.5, "bold": True}),
     ("→ CamelBridge.analyze([\"يذهب\"])   LRU-cached", {"size": 9.5, "color": MUTED})],
    fill=TEAL_L, line=TEAL, space=2)

box(s, 8.35, 1.30, 4.42, 1.35,
    [(".venv-camel  (subprocess)", {"size": 10, "bold": True, "color": AMBER}),
     ("MLEDisambiguator.disambiguate([…])", {"size": 11.5, "bold": True}),
     ("→ dr.analyses → _trim() to 11 fields", {"size": 9.5, "color": MUTED})],
    fill=AMBER_L, line=AMBER, space=2)

a = arrow(s, 7.80, 1.52, 0.45, 0.20, "right", AMBER)
t = s.shapes.add_textbox(Inches(7.65), Inches(1.70), Inches(0.68), Inches(0.24))
text(t, "stdin", size=8.5, color=AMBER, bold=True)
arrow(s, 7.80, 2.18, 0.45, 0.20, "left", AMBER)
t = s.shapes.add_textbox(Inches(7.65), Inches(2.36), Inches(0.68), Inches(0.24))
text(t, "stdout", size=8.5, color=AMBER, bold=True)

box(s, 0.55, 2.85, 12.22, 0.95,
    [('{"root": "ذ.ه.ب",  "pattern": "يَ1ْ2َ3",  "stem": "ذْهَب",  "diac": "يَذْهَب",  "lex": "ذَهَب",  "pos": "verb",',
      {"size": 11.5, "font": MONO, "color": RGBColor(0xE2, 0xE8, 0xF0), "align": PP_ALIGN.LEFT}),
     ('  "prc3": "0", "prc2": "0", "prc1": "0", "prc0": "0", "enc0": "0"}          ← the raw candidate, verbatim',
      {"size": 11.5, "font": MONO, "color": RGBColor(0xE2, 0xE8, 0xF0), "align": PP_ALIGN.LEFT})],
    fill=RGBColor(0x1E, 0x29, 0x3B), line=None, radius=0.06)

# the five post-processing rules
t = s.shapes.add_textbox(Inches(0.55), Inches(3.95), Inches(6.0), Inches(0.3))
text(t, "_dict_to_analysis()  —  five transformations, two of them rejections",
     size=12, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

rules = [
    ("1", "strip root separators  _  .  #", "'ذ.ه.ب'  →  'ذهب'", GREEN),
    ("2", "reject if root < 3 letters", "3 letters → pass", GREEN),
    ("3", "reject NTWS (loanword marker)", "not a loanword → pass", GREEN),
    ("4", "clitic feature tags → Arabic surface", 'all "0"  →  None', GREEN),
    ("5", "normalize_pattern() strips clitic chars", "no clitics → unchanged", GREEN),
]
y = 4.35
for num, what, res, col in rules:
    box(s, 0.55, y, 0.42, 0.46, num, fill=TEAL, line=None, size=12, bold=True, color=WHITE, radius=0.3)
    box(s, 1.05, y, 4.55, 0.46, what, fill=WHITE, line=SLATE_L, size=11,
        align=PP_ALIGN.LEFT, bold=False)
    box(s, 5.68, y, 3.05, 0.46, res, fill=GREEN_L, line=None, size=10.5, color=GREEN, bold=True)
    y += 0.55

box(s, 9.00, 4.35, 3.77, 2.30,
    [("Analysis  (frozen dataclass)", {"size": 11.5, "bold": True, "color": TEAL}),
     ("root      = 'ذهب'", {"size": 11, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("pattern   = 'يَ1ْ2َ3'   ← bare stem", {"size": 11, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("stem      = 'ذْهَب'    ← NOT used", {"size": 11, "font": MONO, "align": PP_ALIGN.LEFT, "color": ROSE}),
     ("surface   = 'يَذْهَب'   ← used", {"size": 11, "font": MONO, "align": PP_ALIGN.LEFT, "color": GREEN}),
     ("prc0-3, enc0 = None", {"size": 11, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("", {"size": 4}),
     ("CAMeL's `stem` drops the present-tense ي —\nreconstruction uses diac minus clitics instead,\nso the conjugation survives.",
      {"size": 9.5, "color": MUTED, "align": PP_ALIGN.LEFT})],
    fill=WHITE, line=TEAL, space=3)

footnote(s, "Rules 2 and 3 are where 51 % of real-corpus word occurrences are lost — see the diagnostics slide.")

# ==========================================================================
# Slide 4 — vocabulary layout and the two hits
# ==========================================================================
s = blank(prs)
title(s, "Vocabulary lookup — deterministic ID ranges", "Emission is all-or-nothing per word: if either token is OOV, the whole word falls back to characters.",
      accent=INDIGO)

segs = [
    ("specials", "0–3", 4, SLATE, 0.97),
    ("LIT", "4–5", 2, SLATE, 0.66),
    ("proclitics", "6–17", 12, TEAL, 1.20),
    ("enclitics", "18–29", 12, TEAL, 1.20),
    ("chars", "30–76", 47, AMBER, 0.97),
    ("digits", "77–96", 20, AMBER, 0.87),
    ("punct", "97–137", 41, AMBER, 0.92),
    ("ROOTS", "138–3283", 3146, INDIGO, 2.95),
    ("PATTERNS", "3284–3783", 500, PURPLE, 1.93),
]
x = 0.55
for name, rng, n, col, w in segs:
    hot = name in ("ROOTS", "PATTERNS")
    box(s, x, 1.45, w, 0.95,
        [(name, {"size": 11.5 if hot else 9.5, "bold": True,
                 "color": WHITE if hot else col}),
         (rng, {"size": 8.5, "color": RGBColor(0xE2,0xE8,0xF0) if hot else MUTED}),
         (f"{n:,}", {"size": 8.5, "color": RGBColor(0xE2,0xE8,0xF0) if hot else MUTED})],
        fill=col if hot else WHITE, line=col, lw=1.6 if hot else 1.0, space=1)
    x += w + 0.07

box(s, 0.55, 2.62, 12.22, 0.42, "3,784 tokens total   ·   max_roots: 10000 never reached (only 3,146 roots occur ≥ 2×)   ·   max_patterns: 500 WAS binding — 6,076 qualified",
    fill=INDIGO_L, line=None, size=11, color=INDIGO, bold=True, radius=0.15)

# the two hits
box(s, 0.55, 3.35, 5.9, 1.35,
    [("[ROOT_ذهب]", {"size": 20, "bold": True, "color": INDIGO}),
     ("id 354      freq 346      source: corpus", {"size": 11, "font": MONO, "color": MUTED}),
     ("examples: الذهبية · مذهب · تذهب · ذهب · الذهبي", {"size": 11, "color": INK})],
    fill=WHITE, line=INDIGO, lw=1.8, space=3)
box(s, 6.85, 3.35, 5.92, 1.35,
    [("[PAT_يَ1ْ2َ3]", {"size": 20, "bold": True, "color": PURPLE}),
     ("id 3298     freq 1899     source: corpus", {"size": 11, "font": MONO, "color": MUTED}),
     ("examples: يَرْغَب · يَشْمَل · يَجْمَع · لِيَسْمَح", {"size": 11, "color": INK})],
    fill=WHITE, line=PURPLE, lw=1.8, space=3)

# the grid claim
box(s, 0.55, 4.95, 12.22, 1.25,
    [("Why a pair beats a subword entry", {"size": 12.5, "bold": True, "color": INK}),
     ("root ذهب appears with 58 distinct patterns   ×   pattern يَ1ْ2َ3 appears with 360 distinct roots",
      {"size": 13, "bold": True, "color": INDIGO}),
     ("ذهب · ذاهب · يذهب · تذهب · مذهب · ذهاب · مذاهب · أذهب · ذهبت …          يجمع · يعلم · يعمل · يقطع · يدفع · يحفظ · يظهر …",
      {"size": 10.5, "color": MUTED})],
    fill=RGBColor(0xF1, 0xF5, 0xF9), line=SLATE_L, space=4)

box(s, 0.55, 6.35, 12.22, 0.52,
    [("418 surface forms are addressable with 2 tokens instead of 418 separate vocabulary entries — that is the compression claim.",
      {"size": 12, "bold": True, "color": WHITE})],
    fill=INDIGO, line=None, radius=0.14)

# ==========================================================================
# Slide 5 — into the model
# ==========================================================================
s = blank(prs)
title(s, "Into the model — a gather, 16 layers, and a tied projection",
      "embedding_type = standard, so LLaMA is untouched apart from the size of one matrix.", accent=PURPLE)

box(s, 0.55, 1.40, 2.55, 0.75, [("[1, 354, 3298, 2]", {"size": 13, "bold": True, "font": MONO, "color": WHITE})],
    fill=INDIGO, line=None)
arrow(s, 3.22, 1.66, 0.22, 0.24, "right", PURPLE)

box(s, 3.60, 1.28, 3.35, 2.35,
    [("embed_tokens", {"size": 13, "bold": True, "color": PURPLE}),
     ("nn.Embedding(3784, 2048)", {"size": 11, "font": MONO}),
     ("7,749,632 parameters · bf16", {"size": 9.5, "color": MUTED}),
     ("", {"size": 4}),
     ("E[354]  ‖v‖ = 1.1228", {"size": 10.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("E[3298] ‖v‖ = 1.1613", {"size": 10.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("", {"size": 4}),
     ("gather 4 rows → [B, 4, 2048]", {"size": 11, "bold": True, "color": PURPLE})],
    fill=PURPLE_L, line=PURPLE, space=2)

box(s, 7.20, 1.28, 2.55, 2.35,
    [("16 decoder layers", {"size": 13, "bold": True, "color": PURPLE}),
     ("", {"size": 4}),
     ("RMSNorm", {"size": 10.5}), ("RoPE θ=500000", {"size": 10.5}),
     ("GQA 32 heads / 8 KV", {"size": 10.5}), ("SwiGLU 2048→8192", {"size": 10.5}),
     ("", {"size": 4}),
     ("[B, 4, 2048]", {"size": 11, "bold": True, "color": PURPLE})],
    fill=WHITE, line=PURPLE, space=1)

box(s, 10.00, 1.28, 2.77, 2.35,
    [("lm_head", {"size": 13, "bold": True, "color": PURPLE}),
     ("tied to embed_tokens", {"size": 10.5, "color": ROSE, "bold": True}),
     ("", {"size": 4}),
     ("2048 → 3,784", {"size": 11, "font": MONO}),
     ("logits [B, 4, 3784]", {"size": 11, "bold": True, "color": PURPLE}),
     ("", {"size": 4}),
     ("absent from named_parameters()\n— training embed_tokens IS\ntraining lm_head",
      {"size": 9, "color": MUTED})],
    fill=PURPLE_L, line=PURPLE, space=2)
arrow(s, 7.00, 2.33, 0.16, 0.22, "right", PURPLE)
arrow(s, 9.80, 2.33, 0.16, 0.22, "right", PURPLE)

# the inherited rows
box(s, 0.55, 3.85, 6.10, 1.75,
    [("What row 354 was before Phase 1", {"size": 12.5, "bold": True, "color": ROSE}),
     ("resize 128,256 → 3,784 is a SHRINK: HuggingFace keeps the first 3,784 pretrained rows, so nothing is reinitialised.",
      {"size": 10, "color": MUTED}),
     ("", {"size": 3}),
     ("[ROOT_ذهب]   inherits LLaMA's  'ot'", {"size": 12, "font": MONO, "color": INK}),
     ("[PAT_يَ1ْ2َ3]  inherits LLaMA's  ' Act'", {"size": 12, "font": MONO, "color": INK}),
     ("[CHAR_ا]      inherits LLaMA's  'E'", {"size": 12, "font": MONO, "color": INK}),
     ("An index collision, not an alignment. Phase 1 exists to drift these rows.", {"size": 10, "color": ROSE, "bold": True})],
    fill=ROSE_L, line=ROSE, space=2)

# real top-k
box(s, 6.95, 3.85, 5.82, 1.75,
    [("Real top-4 after «…الإجابة:»  (Phase-3 checkpoint)", {"size": 12, "bold": True, "color": PURPLE}),
     ("[PAT_1َ2َ3َ]     p = 0.7225      id 3320", {"size": 11.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[PAT_يَ1ْ2ُ3]    p = 0.2070      id 3299", {"size": 11.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[PAT_مَ1ْ2َ3]    p = 0.0080      id 3306", {"size": 11.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[ROOT_نشط]      p = 0.0075      id  330", {"size": 11.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("94 % of the top-8 mass sits on PAT tokens — the model learned the class grammar of the vocabulary.",
      {"size": 10, "color": PURPLE, "bold": True})],
    fill=WHITE, line=PURPLE, space=2)

# training phases strip
phases = [("Phase 1 · embedding alignment", "embed_tokens only, body frozen · Arabic-SQuAD · 1000 steps · loss 1.5933", TEAL),
          ("Phase 2 · warmup", "all params · Arabic-SQuAD · answer-only · 2000 steps · loss 1.6159", AMBER),
          ("Phase 3 · SFT", "all params · TyDiQA-ar + ARCD · 2000 steps · loss 0.2708", GREEN)]
x = 0.55
for name, body, col in phases:
    box(s, x, 5.85, 4.02, 0.95, [(name, {"size": 11.5, "bold": True, "color": col}),
                                 (body, {"size": 9.5, "color": MUTED})],
        fill=WHITE, line=col, space=2)
    x += 4.09

# ==========================================================================
# Slide 6 — decode
# ==========================================================================
s = blank(prs)
title(s, "Decode — one left-to-right pass over four pieces of state",
      "Dispatch is on the token's string prefix, which is why proclitics and enclitics have distinct prefixes.", accent=GREEN)

state = [("out", "finalized words, in order"), ("clitic_prefix", "proclitics waiting for the next content word"),
         ("pending_root", "a ROOT that has not yet met its PAT"), ("lit_buffer", "chars between LIT_BEGIN and LIT_END")]
x = 0.55
for n, d in state:
    box(s, x, 1.40, 3.02, 0.72, [(n, {"size": 12, "bold": True, "font": MONO, "color": GREEN}),
                                 (d, {"size": 9, "color": MUTED})],
        fill=GREEN_L, line=GREEN, space=1)
    x += 3.09

# walk table
hdr = [("token", 2.60), ("branch taken", 5.40), ("state after", 4.22)]
x = 0.55
for h, w in hdr:
    box(s, x, 2.35, w, 0.42, h, fill=SLATE, line=None, size=11, bold=True, color=WHITE, radius=0.05,
        align=PP_ALIGN.LEFT)
    x += w
rows = [
    ("<s>  (1)", "skipped — pad / bos / eos", "—"),
    ("[CLITICP_و]  (6)", "clitic_prefix.append('و')", "prefix = ['و']"),
    ("[ROOT_ذهب]  (354)", "dump_orphan_root(); pending_root = 'ذهب'", "pending = 'ذهب'"),
    ("[PAT_يَ1ْ2َ3ُونَ]  (3423)", "_reconstruct(354, 3423) → 'يذهبون'; flush prepends 'و'", "out = ['ويذهبون']"),
    ("</s>  (2)", "skipped", "—"),
]
y = 2.79
for i, (a_, b_, c_) in enumerate(rows):
    fill = WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF5, 0xF9)
    hot = i == 3
    x = 0.55
    for val, w in zip((a_, b_, c_), (2.60, 5.40, 4.22)):
        box(s, x, y, w, 0.48, val, fill=GREEN_L if hot else fill, line=SLATE_L, lw=0.75,
            size=10.5, bold=hot, color=GREEN if hot else INK, align=PP_ALIGN.LEFT, radius=0.02)
        x += w
    y += 0.48

box(s, 0.55, 5.30, 5.9, 0.55, [("result: «ويذهبون» — exact round trip", {"size": 13, "bold": True, "color": WHITE})],
    fill=GREEN, line=None, radius=0.12)

# three tiers
t = s.shapes.add_textbox(Inches(6.85), Inches(5.22), Inches(6.0), Inches(0.3))
text(t, "_reconstruct() — three tiers", size=12, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
tiers = [("TIER 1", "lookup table  (60,248 entries)", "hits ~99 % · O(1)", GREEN, GREEN_L, True),
         ("TIER 2", "CAMeL Generator via the bridge", "naive-fill → re-analyse → return stem", AMBER, WHITE, False),
         ("TIER 3", "naive slot substitution", "wrong on weak roots: قَوَلَ not قال", ROSE, WHITE, False)]
y = 5.55
for name, what, note, col, fill, hot in tiers:
    box(s, 6.85, y, 1.05, 0.42, name, fill=col, line=None, size=10, bold=True, color=WHITE, radius=0.2)
    box(s, 8.00, y, 4.77, 0.42, [(f"{what}   —   {note}", {"size": 10, "align": PP_ALIGN.LEFT,
        "color": col if hot else INK, "bold": hot})], fill=fill, line=col, lw=1.4 if hot else 0.9)
    y += 0.48

footnote(s, "For يذهب tier 1 hits and tiers 2–3 never run — the expected case, since the LM was trained on exactly this pair distribution.")

# ==========================================================================
# Slide 7 — three branches
# ==========================================================================
s = blank(prs)
title(s, "Same machinery, three outcomes", "Which branch a word takes decides whether AraRooPat behaves like a morphological tokenizer or a character tokenizer.",
      accent=AMBER)

cols = [
    ("يذهب", TEAL, TEAL_L, "root + pattern",
     ["CAMeL: root 'ذ.ه.ب' ✓", "pattern يَ1ْ2َ3 in vocab ✓", "no clitics"],
     ["<s>", "[ROOT_ذهب]", "[PAT_يَ1ْ2َ3]", "</s>"],
     "[1, 354, 3298, 2]", "2 content tokens"),
    ("ويذهبون", INDIGO, INDIGO_L, "clitic split off",
     ["prc2 = 'wa_part' → 'و'", "pattern وَيَ1ْ2َ3ُونَ", "→ strip 'و' → يَ1ْ2َ3ُونَ"],
     ["<s>", "[CLITICP_و]", "[ROOT_ذهب]", "[PAT_يَ1ْ2َ3ُونَ]", "</s>"],
     "[1, 6, 354, 3423, 2]", "3 content tokens · same root"),
    ("الولد", ROSE, ROSE_L, "LIT fallback",
     ["CAMeL: root '#.ل.د'", "strip # → 'لد'  (2 letters)", "REJECTED by the length rule"],
     ["<s>", "[LIT_BEGIN]", "[CHAR_ا][CHAR_ل][CHAR_و][CHAR_ل][CHAR_د]", "[LIT_END]", "</s>"],
     "[1, 4, 36, 59, 63, 59, 44, 5, 2]", "7 tokens for a 5-character word"),
]
x = 0.55
for word, col, light, kind, notes, toks, ids, verdict in cols:
    box(s, x, 1.40, 4.02, 0.80, [(word, {"size": 26, "bold": True, "color": WHITE})], fill=col, line=None)
    box(s, x, 2.28, 4.02, 0.36, kind, fill=light, line=None, size=11, bold=True, color=col, radius=0.18)
    y = 2.76
    for n in notes:
        box(s, x, y, 4.02, 0.38, n, fill=WHITE, line=SLATE_L, lw=0.8, size=10,
            align=PP_ALIGN.LEFT)
        y += 0.42
    y = 4.14
    for tk in toks:
        box(s, x, y, 4.02, 0.36, tk, fill=light, line=col, lw=0.9, size=10, bold=True, color=col)
        y += 0.41
    box(s, x, 6.24, 4.02, 0.42, ids, fill=SLATE, line=None, size=10.5, bold=True, color=WHITE,
        font=MONO, radius=0.08)
    box(s, x, 6.70, 4.02, 0.28, verdict, fill=WHITE, line=None, size=10, color=col, bold=True)
    x += 4.09

# ==========================================================================
# Slide 8 — pattern notation
# ==========================================================================
s = blank(prs)
title(s, "Why a pattern contains 1 2 3 — and where real non-Arabic characters do leak in",
      "The digits are CAMeL's positional slots, not text. The classical Arabic wazn writes the same thing with ف ع ل.",
      accent=PURPLE)

box(s, 0.55, 1.40, 12.22, 0.62,
    [("A pattern is a template with numbered holes:   digit n = the n-th root consonant   ·   every other character is literal template material",
      {"size": 12.5, "bold": True, "color": INK})],
    fill=PURPLE_L, line=None, radius=0.12)

hdr = [("CAMeL slot notation", 3.30), ("classical wazn (ف ع ل)", 3.30), ("root", 1.85), ("surface", 3.77)]
x = 0.55
for h, w in hdr:
    box(s, x, 2.20, w, 0.45, h, fill=SLATE, line=None, size=11, bold=True, color=WHITE, radius=0.05)
    x += w
demo = [("يَ1ْ2َ3", "ذهب", "يذهب"), ("1َ2َ3َ", "ذهب", "ذهب"), ("مَ1ْ2َ3", "ذهب", "مذهب"),
        ("1ا2ِ3", "ذهب", "ذاهب"), ("مَ1ا2ِ3", "ذهب", "مذاهب"), ("1َ2ِي3", "كتب", "كتيب")]
y = 2.68
for pat, root, surf in demo:
    x = 0.55
    for val, w, col, mono in ((pat, 3.30, PURPLE, False), (to_wazn(pat), 3.30, GREEN, False),
                              (root, 1.85, INDIGO, False), (surf, 3.77, INK, False)):
        box(s, x, y, w, 0.46, val, fill=WHITE, line=SLATE_L, lw=0.8, size=13, bold=True, color=col)
        x += w
    y += 0.51

box(s, 0.55, 5.85, 6.05, 1.02,
    [("Verified: 0 of 500 pattern tokens contain a non-Arabic letter", {"size": 11.5, "bold": True, "color": GREEN}),
     ("Every [PAT_*] body is Arabic letters + diacritics + slot digits 1–4 only.", {"size": 10, "color": MUTED})],
    fill=GREEN_L, line=GREEN, space=2)

box(s, 6.72, 5.85, 6.05, 1.02,
    [("But 6 vocabulary entries DO leak Latin text", {"size": 11.5, "bold": True, "color": ROSE}),
     ("[ROOT_FOREIGN] (freq 61) · [ROOT_Uٌٍ] (3) · [CLITICP_la_emph] (12,676) ·\n[CLITICP_la_rc] (1,904) · [CLITICP_>a_ques] (8) · [CLITICE_mA_sub] (6)",
      {"size": 9.5, "color": MUTED})],
    fill=ROSE_L, line=ROSE, space=2)

footnote(s, "Untranslated CAMeL feature tags + a FOREIGN root marker the NTWS guard does not catch. Fixed in araroopat_backend.py — takes effect on the next tokenizer training run.")


# ==========================================================================
# Slide 9 — transformer internals, with the real tensor handed between boxes
# ==========================================================================
s = blank(prs)
title(s, "Inside the transformer — what each box hands to the next",
      "Real forward pass of «يذهب» (4 tokens, fp32). ‖·‖ values are the L2 norm at position 2 = [PAT_يَ1ْ2َ3].",
      accent=PURPLE)

SX, SW = 0.55, 4.55          # left stack
IX, IW = 0.72, 4.21          # boxes inside the layer container


def stack(y, h, lines, fill=WHITE, line=PURPLE, x=SX, w=SW, **kw):
    return box(s, x, y, w, h, lines, fill=fill, line=line, **kw)


def passed(txt):
    return (txt, {"size": 8.5, "color": PURPLE, "bold": True, "font": MONO})


stack(1.22, 0.40, [("input_ids  [1, 4]", {"size": 11, "bold": True, "color": WHITE}),
                   ("[ 1, 354, 3298, 2 ]", {"size": 8.5, "font": MONO, "color": RGBColor(0xC7, 0xD2, 0xFE)})],
      fill=INDIGO, line=None, space=0)
arrow(s, 2.75, 1.64, 0.16, 0.12, "down", PURPLE)

stack(1.80, 0.44, [("embed_tokens — Embedding(3784, 2048)", {"size": 10.5, "bold": True}),
                   passed("out  [1, 4, 2048]      ‖·‖ = 1.161")], space=0)
arrow(s, 2.75, 2.26, 0.16, 0.12, "down", PURPLE)

# ---- the repeated decoder layer -------------------------------------------
box(s, SX, 2.42, SW, 3.42, fill=RGBColor(0xF5, 0xF3, 0xFF), line=PURPLE, lw=1.6)
box(s, 0.60, 2.45, 4.45, 0.30, "decoder layer  ×  16   (LlamaDecoderLayer)",
    fill=PURPLE, line=None, size=10, bold=True, color=WHITE, radius=0.12)

inner = [
    ("input_layernorm — RMSNorm(2048)", "out  [1, 4, 2048]      ‖·‖ = 9.953", WHITE),
    ("q_proj 2048→2048 (32 h)  ·  k_proj →512 (8 h)  ·  v_proj →512",
     "q ‖·‖ = 76.4    k ‖·‖ = 41.4    v ‖·‖ = 2.53", WHITE),
    ("RoPE — rotate q and k by position   ← position enters HERE",
     "cos, sin  [1, 4, 64]   ·   0 trainable parameters", AMBER_L),
    ("softmax(Q Kᵀ / √64) V   causal  →  o_proj 2048→2048",
     "attn out  [1, 4, 2048]      ‖·‖ = 1.923", WHITE),
    ("⊕  residual:  hidden = hidden + attn_out", "‖·‖ = 1.161 ⊕ 1.923", GREEN_L),
    ("post_attention_layernorm — RMSNorm", "out  ‖·‖ = 9.611", WHITE),
    ("MLP  down( SiLU(gate(x)) ⊙ up(x) )   2048→8192→2048",
     "mlp out  [1, 4, 2048]      ‖·‖ = 10.540", WHITE),
    ("⊕  residual:  hidden = hidden + mlp_out", "layer output  ‖·‖ = 11.095", GREEN_L),
]
y = 2.80
for what, val, fill in inner:
    box(s, IX, y, IW, 0.34,
        [(what, {"size": 8.8, "bold": True, "align": PP_ALIGN.LEFT}),
         (val, {"size": 7.6, "font": MONO, "align": PP_ALIGN.LEFT,
                "color": AMBER if fill is AMBER_L else (GREEN if fill is GREEN_L else PURPLE)})],
        fill=fill, line=SLATE_L, lw=0.7, radius=0.06, space=0)
    y += 0.38

arrow(s, 2.75, 5.86, 0.16, 0.12, "down", PURPLE)
stack(6.02, 0.40, [("model.norm — RMSNorm(2048)", {"size": 10.5, "bold": True}),
                   passed("out  [1, 4, 2048]      ‖·‖ = 46.817")], space=0)
arrow(s, 2.75, 6.46, 0.16, 0.12, "down", PURPLE)
stack(6.62, 0.50, [("lm_head — tied Wᵀ, 2048 → 3,784", {"size": 10.5, "bold": True, "color": WHITE}),
                   ("logits [1, 4, 3784]  →  log_softmax  →  argmax", {"size": 8.5, "font": MONO,
                    "color": RGBColor(0xDD, 0xD6, 0xFE)})],
      fill=PURPLE, line=None, space=0)

# ---- panel A: token embedding + position ----------------------------------
box(s, 5.30, 1.22, 7.48, 2.62, fill=WHITE, line=SLATE_L)
t = s.shapes.add_textbox(Inches(5.45), Inches(1.26), Inches(7.2), Inches(0.30))
text(t, "«token + position» — where each one actually enters", size=12.5, bold=True,
     color=INK, align=PP_ALIGN.LEFT)

box(s, 5.45, 1.62, 3.55, 1.86,
    [("1 · TOKEN  — a lookup, nothing more", {"size": 10, "bold": True, "color": INDIGO}),
     ("E[354]   ROOT_ذهب", {"size": 9, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("  [-0.0092, -0.0187, -0.0160, …]  ‖·‖=1.1228", {"size": 8.2, "font": MONO, "align": PP_ALIGN.LEFT, "color": MUTED}),
     ("E[3298]  PAT_يَ1ْ2َ3", {"size": 9, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("  [-0.0018, +0.0072, +0.0081, …]  ‖·‖=1.1613", {"size": 8.2, "font": MONO, "align": PP_ALIGN.LEFT, "color": MUTED}),
     ("hidden_states[0] IS this gather —", {"size": 9, "color": INDIGO, "bold": True}),
     ("nothing is summed into it.", {"size": 9, "color": INDIGO, "bold": True})],
    fill=INDIGO_L, line=INDIGO, space=1)

box(s, 9.15, 1.62, 3.48, 1.86,
    [("2 · POSITION — no vector is added", {"size": 10, "bold": True, "color": AMBER}),
     ("LLaMA-3.2 has no wpe / learned pos table.", {"size": 8.6, "color": MUTED}),
     ("inv_freq = 1 / 500000^(2i/64)   ·  32 values", {"size": 8.4, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("  [1.0, 0.6636, 0.4404, 0.2922, …]", {"size": 8.4, "font": MONO, "align": PP_ALIGN.LEFT, "color": MUTED}),
     ("pos0  cos [1.000, 1.000]   sin [0.000, 0.000]", {"size": 8.4, "font": MONO, "align": PP_ALIGN.LEFT, "color": MUTED}),
     ("pos2  cos [-0.416, 0.241]  sin [0.909, 0.971]", {"size": 8.4, "font": MONO, "align": PP_ALIGN.LEFT, "color": MUTED}),
     ("q′ = q⊙cos + rotate_half(q)⊙sin   (q, k only)", {"size": 9, "bold": True, "color": AMBER})],
    fill=AMBER_L, line=AMBER, space=1)

# ---- panel B: real attention map ------------------------------------------
box(s, 5.30, 3.98, 3.55, 1.92, fill=WHITE, line=SLATE_L)
t = s.shapes.add_textbox(Inches(5.42), Inches(4.02), Inches(3.3), Inches(0.26))
text(t, "attention · layer 15, head 0", size=10, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
cols = ["<s>", "ROOT", "PAT", "</s>"]
box(s, 5.42, 4.33, 1.00, 0.24, "query ↓ key →", fill=WHITE, line=None, size=7.2, color=MUTED)
for k, c in enumerate(cols):
    box(s, 6.42 + k * 0.58, 4.33, 0.56, 0.24, c, fill=SLATE_L, line=None, size=7.5,
        bold=True, color=INK, radius=0.1)
attn = [(0, [1.000, 0, 0, 0]), (1, [0.517, 0.483, 0, 0]),
        (2, [0.000, 0.000, 1.000, 0]), (3, [0.000, 0.000, 0.110, 0.890])]
for r, (i, row) in enumerate(attn):
    box(s, 5.42, 4.59 + r * 0.26, 1.00, 0.24, cols[i], fill=SLATE_L, line=None,
        size=7.5, bold=True, color=INK, radius=0.1)
    for k, v in enumerate(row):
        hot = v > 0.4
        box(s, 6.42 + k * 0.58, 4.59 + r * 0.26, 0.56, 0.24,
            ("—" if v == 0 and k > i else f"{v:.2f}"),
            fill=PURPLE if hot else (PURPLE_L if v > 0 else RGBColor(0xF8, 0xFA, 0xFC)),
            line=None, size=7.8, bold=hot,
            color=WHITE if hot else (PURPLE if v > 0 else RGBColor(0xCB, 0xD5, 0xE1)), radius=0.08)

# ---- panel C: the residual stream is literally a sum -----------------------
box(s, 9.05, 3.98, 3.73, 1.92, fill=WHITE, line=SLATE_L)
t = s.shapes.add_textbox(Inches(9.17), Inches(4.02), Inches(3.5), Inches(0.26))
text(t, "the residual stream is literally a sum", size=10, bold=True, color=GREEN,
     align=PP_ALIGN.LEFT)
box(s, 9.17, 4.33, 3.49, 1.45,
    [("layer0_out = embed + attn_out + mlp_out", {"size": 9.5, "bold": True, "font": MONO,
      "align": PP_ALIGN.LEFT, "color": GREEN}),
     ("", {"size": 3}),
     ("embed      ‖·‖ =  1.1613", {"size": 9, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("attn_out   ‖·‖ =  1.9228", {"size": 9, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("mlp_out    ‖·‖ = 10.5403", {"size": 9, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("layer0_out ‖·‖ = 11.0952", {"size": 9, "font": MONO, "align": PP_ALIGN.LEFT, "bold": True}),
     ("torch.allclose(l0, e + attn + mlp) → True", {"size": 8.6, "color": GREEN, "bold": True})],
    fill=GREEN_L, line=GREEN, space=0)

# ---- panel D: the value at every position, every depth ---------------------
t = s.shapes.add_textbox(Inches(5.30), Inches(5.96), Inches(7.4), Inches(0.24))
text(t, "‖hidden‖ at each position, by depth", size=10, bold=True, color=INK, align=PP_ALIGN.LEFT)
head = ["", "pos0  <s>", "pos1  ROOT", "pos2  PAT", "pos3  </s>"]
rows = [("embed_tokens", "0.932", "1.123", "1.161", "0.943"),
        ("layer 0 out", "13.856", "7.760", "11.095", "8.898"),
        ("layer 15 out", "1622.452", "312.784", "203.566", "173.416"),
        ("after model.norm", "68.743", "89.045", "46.817", "45.665")]
wcol = [2.00, 1.37, 1.37, 1.37, 1.37]
x = 5.30
for h_, w_ in zip(head, wcol):
    box(s, x, 6.22, w_, 0.22, h_, fill=SLATE, line=None, size=7.8, bold=True, color=WHITE, radius=0.04)
    x += w_
for r, row in enumerate(rows):
    x = 5.30
    for c, (val, w_) in enumerate(zip(row, wcol)):
        hot = c == 3
        box(s, x, 6.44 + r * 0.22, w_, 0.22, val,
            fill=PURPLE_L if hot else (WHITE if r % 2 == 0 else RGBColor(0xF1, 0xF5, 0xF9)),
            line=SLATE_L, lw=0.5, size=7.8, bold=(c == 0 or hot),
            color=PURPLE if hot else INK,
            align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER, radius=0.02)
        x += w_

# reorder: the transformer-internals slide belongs right after the model slide
_ids = prs.slides._sldIdLst
_last = list(_ids)[-1]
_ids.remove(_last)
_ids.insert(5, _last)

for p in audit(prs):
    print("  !", p)

prs.save(OUT)
print("wrote", OUT, OUT.stat().st_size, "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
