"""Shared presentation kit for the docs/*.pptx generators.

Palette, box/arrow/text primitives, and a geometry audit that refuses to let a
shape leave the slide or two text boxes overlap. Used by
``araroopat_roundtrip_deck.py`` and ``araroopat_mechanism_deck.py``.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Palette — one hue per pipeline stage, reused on every slide.
# --------------------------------------------------------------------------
INK        = RGBColor(0x0F, 0x17, 0x2A)   # near-black text
MUTED      = RGBColor(0x64, 0x74, 0x8B)   # secondary text
PAGE       = RGBColor(0xF8, 0xFA, 0xFC)   # slide background
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

SLATE      = RGBColor(0x1E, 0x29, 0x3B)   # input / output
SLATE_L    = RGBColor(0xE2, 0xE8, 0xF0)
TEAL       = RGBColor(0x0F, 0x76, 0x6E)   # tokenizer front-end
TEAL_L     = RGBColor(0xCC, 0xFB, 0xF1)
AMBER      = RGBColor(0xB4, 0x53, 0x09)   # CAMeL subprocess
AMBER_L    = RGBColor(0xFE, 0xF3, 0xC7)
INDIGO     = RGBColor(0x43, 0x38, 0xCA)   # vocabulary / ids
INDIGO_L   = RGBColor(0xE0, 0xE7, 0xFF)
PURPLE     = RGBColor(0x6D, 0x28, 0xD9)   # model
PURPLE_L   = RGBColor(0xED, 0xE9, 0xFE)
GREEN      = RGBColor(0x15, 0x80, 0x3D)   # decode
GREEN_L    = RGBColor(0xDC, 0xFC, 0xE7)
ROSE       = RGBColor(0xBE, 0x12, 0x3C)   # LIT fallback / problems
ROSE_L     = RGBColor(0xFF, 0xE4, 0xE6)

AR = "Arial"          # has Arabic coverage on Windows + macOS
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)


# --------------------------------------------------------------------------
# Low-level helpers
# --------------------------------------------------------------------------

def _set_cs_font(run, typeface):
    """python-pptx only sets the latin typeface; Arabic uses the complex-script
    slot. Without this PowerPoint falls back to Times-ish for Arabic runs."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:cs", "a:ea"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", typeface)


def _rtl(paragraph, on=True):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1" if on else "0")


def text(shape, lines, size=12, color=INK, bold=False, font=AR,
         align=PP_ALIGN.CENTER, rtl=False, space=0):
    """lines: str | list[str] | list[(str, dict overrides)]"""
    tf = shape.text_frame if hasattr(shape, "text_frame") else shape
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        opts = {}
        if isinstance(ln, tuple):
            ln, opts = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.space_after = Pt(opts.get("space", space))
        _rtl(p, opts.get("rtl", rtl))
        r = p.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.color.rgb = opts.get("color", color)
        f.name = opts.get("font", font)
        _set_cs_font(r, opts.get("font", font))
    return tf


def box(slide, x, y, w, h, lines=None, fill=WHITE, line=SLATE_L, lw=1.25,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10, **kw):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if lines is not None:
        text(s, lines, **kw)
    return s


def arrow(slide, x, y, w, h, direction="right", color=MUTED):
    shp = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
           "left": MSO_SHAPE.LEFT_ARROW, "up": MSO_SHAPE.UP_ARROW}[direction]
    s = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAGE
    return s


def title(slide, main, sub=None, accent=TEAL):
    box(slide, 0, 0, 13.333, 0.06, fill=accent, line=None,
        shape=MSO_SHAPE.RECTANGLE)
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.62))
    text(t, main, size=26, bold=True, align=PP_ALIGN.LEFT, color=INK)
    if sub:
        u = slide.shapes.add_textbox(Inches(0.58), Inches(0.80), Inches(12.2), Inches(0.34))
        text(u, sub, size=12.5, align=PP_ALIGN.LEFT, color=MUTED)


def footnote(slide, txt):
    t = slide.shapes.add_textbox(Inches(0.55), Inches(6.98), Inches(12.2), Inches(0.32))
    text(t, txt, size=9.5, align=PP_ALIGN.LEFT, color=MUTED)


def chip(slide, x, y, w, h, label, fill, edge, size=11.5, bold=True, color=None):
    return box(slide, x, y, w, h, label, fill=fill, line=edge, lw=1.1, radius=0.22,
               size=size, bold=bold, color=color or edge)


# --------------------------------------------------------------------------
# Arabic wazn rendering: CAMeL slot digits -> the classical ف ع ل template
# --------------------------------------------------------------------------
WAZN = {"1": "ف", "2": "ع", "3": "ل", "4": "ل"}


def to_wazn(pattern):
    return "".join(WAZN.get(c, c) for c in pattern)



# --------------------------------------------------------------------------
# Geometry self-check: nothing may leave the slide or collide with a sibling
# --------------------------------------------------------------------------

def audit(prs):
    problems = []
    for i, sl in enumerate(prs.slides, 1):
        rects = []
        for sh in sl.shapes:
            if sh.left is None:
                continue
            r = (sh.left, sh.top, sh.left + sh.width, sh.top + sh.height,
                 sh.text_frame.text[:24].replace("\n", " ") if sh.has_text_frame else "")
            if r[2] > prs.slide_width or r[3] > prs.slide_height or r[0] < 0 or r[1] < 0:
                problems.append(f"slide {i}: out of bounds {r[4]!r}")
            if sh.has_text_frame and sh.text_frame.text:
                rects.append(r)
        for a in range(len(rects)):
            for b in range(a + 1, len(rects)):
                A, B = rects[a], rects[b]
                ox = min(A[2], B[2]) - max(A[0], B[0])
                oy = min(A[3], B[3]) - max(A[1], B[1])
                if ox > Inches(0.12) and oy > Inches(0.12):
                    problems.append(f"slide {i}: text overlap {A[4]!r} / {B[4]!r}")
    return problems


def source_ref(path: str, symbol: str) -> str:
    """Return ``basename.py:LINE`` for a def/class, resolved at build time.

    Slides quote line numbers so they can be jumped to while reading the code;
    resolving them here means regenerating the deck re-syncs them instead of
    letting them rot.
    """
    from pathlib import Path as _P
    p = _P(path)
    for i, line in enumerate(p.read_text().splitlines(), 1):
        st = line.strip()
        if st.startswith((f"def {symbol}(", f"class {symbol}(", f"class {symbol}:")):
            return f"{p.name}:{i}"
    return p.name
