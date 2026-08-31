"""Generate docs/araroopat_mechanism.pptx — how the AraRooPat tokenizer works.

A companion to araroopat_roundtrip_deck.py: that deck traces one word end to
end through the model, this one explains the *mechanism* — every path a word
can take, and which function decides. Function references carry line numbers
resolved at build time (see deck_kit.source_ref), so regenerating re-syncs
them against the code.

All examples are real output of outputs/tokenizers/araroopat_hashfix.

    .venv/bin/python docs/araroopat_mechanism_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from deck_kit import *  # noqa: F403
from deck_kit import (W, H, arrow, audit, blank, box, footnote, source_ref, text, title)

OUT = Path(__file__).with_name("araroopat_mechanism.pptx")
REPO = Path(__file__).resolve().parent.parent
TOK = REPO / "src/arabic_eval/tokenizers/araroopat.py"
BACK = REPO / "src/arabic_eval/tokenizers/araroopat_backend.py"
BRIDGE = REPO / "src/arabic_eval/tokenizers/araroopat_bridge.py"
SERVER = REPO / "src/arabic_eval/tools/araroopat_camel_server.py"


def fn(path, name):
    return source_ref(str(path), name)


prs = Presentation()
prs.slide_width, prs.slide_height = W, H

CODE = {"size": 9, "font": MONO, "color": MUTED, "align": PP_ALIGN.LEFT}
BODY = {"size": 10.5, "align": PP_ALIGN.LEFT}


# ==========================================================================
# 1 — title + the four layers
# ==========================================================================
s = blank(prs)
box(s, 0, 0, 13.333, 7.5, fill=RGBColor(0x0F, 0x17, 0x2A), line=None, shape=MSO_SHAPE.RECTANGLE)
box(s, 0, 0, 13.333, 0.10, fill=TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)

t = s.shapes.add_textbox(Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.85))
text(t, "AraRooPat — how the mechanism works", size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
t = s.shapes.add_textbox(Inches(0.92), Inches(2.00), Inches(11.5), Inches(0.5))
text(t, "every word takes exactly one of six paths  ·  each slide names the function that decides",
     size=15, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)

t = s.shapes.add_textbox(Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.3))
text(t, "FOUR LAYERS, TWO PROCESSES", size=10, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

layers = [
    ("tokenizers/araroopat.py", "encode · decode · vocab construction\nthe two state machines live here", TEAL, "main .venv"),
    ("tokenizers/araroopat_backend.py", "Analysis dataclass · root & pattern\nnormalization · clitic tables", INDIGO, "main .venv"),
    ("tokenizers/araroopat_bridge.py", "NDJSON client · spawns the server\nper-request timeouts · fail-loud", PURPLE, "main .venv"),
    ("tools/araroopat_camel_server.py", "imports camel_tools · MLE\ndisambiguation · generation", AMBER, ".venv-camel"),
]
x = 0.9
for name, body, col, env in layers:
    box(s, x, 3.12, 2.85, 1.62,
        [(name, {"size": 10.5, "bold": True, "color": col, "font": MONO}),
         ("", {"size": 3}),
         (body, {"size": 9.5, "color": RGBColor(0xCB, 0xD5, 0xE1)}),
         ("", {"size": 3}),
         (env, {"size": 8.5, "color": col, "bold": True})],
        fill=RGBColor(0x1B, 0x25, 0x3B), line=col, space=0)
    x += 3.0

t = s.shapes.add_textbox(Inches(0.9), Inches(5.05), Inches(11.5), Inches(0.9))
text(t, ["camel-tools pins numpy<2 and transformers<4.54, which breaks lighteval — so it lives in its "
         "own virtualenv and is reached over stdin/stdout NDJSON. There is no in-process import and no "
         "degraded mode: a missing .venv-camel raises CamelBridgeError.",
         ""],
     size=11, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT, space=4)

box(s, 0.9, 5.95, 11.53, 0.72,
    [("الكتاب  →  [CLITICP_ال] [ROOT_كتب] [PAT_1ِ2ا3ِ]  →  الكتاب",
      {"size": 17, "bold": True, "color": RGBColor(0x2D, 0xD4, 0xBF)})],
    fill=RGBColor(0x1B, 0x25, 0x3B), line=None)


# ==========================================================================
# 2 — the decision tree: every path a word can take
# ==========================================================================
s = blank(prs)
title(s, "Every path a word can take", "One decision tree. Each leaf is a different token shape; the colour follows through the whole deck.")

AX, AW = 0.55, 4.60
IX = 5.35

def spine(y, h, lines, fill=WHITE, line=TEAL, **kw):
    return box(s, AX, y, AW, h, lines, fill=fill, line=line, **kw)

spine(1.28, 0.46, [("whitespace word", {"size": 11.5, "bold": True, "color": WHITE})],
      fill=SLATE, line=None)
arrow(s, 2.77, 1.78, 0.16, 0.13, "down", TEAL)
spine(1.95, 0.56,
      [(f"_encode_word()   {fn(TOK,'_encode_word')}", {"size": 9.5, "bold": True, "font": MONO}),
       ("cut into runs by _classify_char(): alpha · digit · punct · other",
        {"size": 8.8, "color": MUTED})], space=0)
arrow(s, 2.77, 2.56, 0.16, 0.13, "down", TEAL)
spine(2.73, 0.50,
      [("run class = alpha", {"size": 10.5, "bold": True, "color": TEAL}),
       (f"_emit_alpha()   {fn(TOK,'_emit_alpha')}", {"size": 8.5, "font": MONO, "color": MUTED})],
      fill=TEAL_L, space=0)
arrow(s, 2.77, 3.28, 0.16, 0.13, "down", TEAL)
spine(3.45, 0.56,
      [("MorphAnalyzer.analyze() → bridge → CAMeL", {"size": 10, "bold": True}),
       ("_first_valid() walks candidates in MLE score order", {"size": 8.5, "color": MUTED})], space=0)
arrow(s, 2.77, 4.06, 0.16, 0.13, "down", TEAL)

spine(4.23, 0.58,
      [("does _dict_to_analysis() accept it?", {"size": 10.5, "bold": True, "color": AMBER}),
       (f"{fn(BACK,'_dict_to_analysis')}   — four gates", {"size": 8.5, "font": MONO, "color": MUTED})],
      fill=AMBER_L, line=AMBER, space=0)
arrow(s, 2.77, 4.86, 0.16, 0.13, "down", TEAL)

spine(5.03, 0.58,
      [("are [ROOT_x] AND [PAT_y] both in vocab?", {"size": 10.5, "bold": True, "color": INDIGO}),
       ("all-or-nothing: one miss discards the other token too", {"size": 8.5, "color": MUTED})],
      fill=INDIGO_L, line=INDIGO, space=0)
arrow(s, 2.77, 5.66, 0.16, 0.13, "down", GREEN)

spine(5.83, 0.86,
      [("① CLITICP* · ROOT · PAT · CLITICE*", {"size": 12, "bold": True, "color": WHITE}),
       ("والكتاب → [CLITICP_و] [CLITICP_ال] [ROOT_كتب] [PAT_1ُ2ّا3ِ]",
        {"size": 9, "color": RGBColor(0xDC, 0xFC, 0xE7)})],
      fill=GREEN, line=None, space=1)

# the two "no" exits
for y, lbl in ((4.23, "no — CAMeL cannot decompose it"), (5.03, "no — out of vocabulary")):
    arrow(s, AX + AW + 0.04, y + 0.22, 0.22, 0.14, "right", ROSE)
    box(s, IX, y, 3.05, 0.58, [("no", {"size": 10, "bold": True, "color": ROSE}),
                               (lbl, {"size": 8.2, "color": MUTED})],
        fill=ROSE_L, line=ROSE, space=0)
arrow(s, 6.85, 5.66, 0.16, 0.13, "down", ROSE)
box(s, IX, 5.83, 3.05, 0.86,
    [("② _emit_lit()", {"size": 11.5, "bold": True, "color": WHITE, "font": MONO}),
     ("[LIT_BEGIN] [CHAR_*]… [LIT_END]", {"size": 8.5, "color": RGBColor(0xFF, 0xE4, 0xE6)})],
    fill=ROSE, line=None, space=1)

# the three non-alpha lanes
lanes = [
    ("run class = digit", "_emit_atom() — one token per character", "③  2024 → [DIGIT_2][DIGIT_0][DIGIT_2][DIGIT_4]", AMBER, AMBER_L),
    ("run class = punct", "_emit_atom() — one token per character", "④  ! → [PUNCT_!]", AMBER, AMBER_L),
    ("run class = other", "Latin · ى · ٱ · emoji — never analysed", "⑤  <unk>", SLATE, SLATE_L),
]
y = 2.73
for head, mid, out, col, light in lanes:
    box(s, 8.75, y, 4.03, 0.44, head, fill=light, line=col, size=10, bold=True, color=col)
    box(s, 8.75, y + 0.48, 4.03, 0.34, mid, fill=WHITE, line=SLATE_L, lw=0.7, size=8.5,
        color=MUTED, align=PP_ALIGN.LEFT)
    box(s, 8.75, y + 0.86, 4.03, 0.40, out, fill=col, line=None, size=9.5, bold=True, color=WHITE)
    arrow(s, 10.68, y + 0.44, 0.14, 0.04, "down", col)
    y += 1.44

footnote(s, "⑥ A clitic surface that is not in the vocabulary for its kind also falls to _emit_lit() — "
            "see _emit_clitic(). Six leaves in total; the whole word takes one path, never a mixture.")


# ==========================================================================
# 3 — path ①: the root + pattern emission
# ==========================================================================
s = blank(prs)
title(s, "Path ① — root + pattern", "The design case: the word becomes two discrete symbols, and the model never sees the surface string.",
      accent=GREEN)

box(s, 0.55, 1.30, 12.22, 0.52,
    [("emission order is fixed — outermost proclitic first, enclitic last:   "
      "prc3 · prc2 · prc1 · prc0 · ROOT · PAT · enc0",
      {"size": 12, "bold": True, "color": GREEN})],
    fill=GREEN_L, line=None, radius=0.12)

rows = [
    ("كتب", "[ROOT_كتب] [PAT_1َ2َ3َ]", "sound triliteral, no clitics"),
    ("يذهب", "[ROOT_ذهب] [PAT_يَ1ْ2َ3]", "the present-tense ي is inside the PATTERN, not a clitic"),
    ("يذهبون", "[ROOT_ذهب] [PAT_يَ1ْ2َ3ُونَ]", "same root, inflection rides the pattern"),
    ("المدرسة", "[CLITICP_ال] [ROOT_درس] [PAT_مَ1ْ2َ3َةِ]", "article split off as its own token"),
    ("وكتابه", "[CLITICP_و] [ROOT_كتب] [PAT_1ِ2ا3] [CLITICE_ه]", "proclitic + enclitic around the stem"),
    ("قرأته", "[ROOT_قر#] [PAT_1َ2َأَت] [CLITICE_ه]", "masked radical + enclitic"),
]
y = 1.98
for w, toks, note in rows:
    box(s, 0.55, y, 1.55, 0.50, w, fill=SLATE, line=None, size=15, bold=True, color=WHITE, radius=0.08)
    box(s, 2.20, y, 6.35, 0.50, toks, fill=GREEN_L, line=GREEN, lw=0.9, size=10.5, bold=True, color=GREEN)
    box(s, 8.65, y, 4.13, 0.50, note, fill=WHITE, line=SLATE_L, lw=0.7, size=9.5,
        color=MUTED, align=PP_ALIGN.LEFT)
    y += 0.56

t = s.shapes.add_textbox(Inches(0.55), Inches(5.45), Inches(6.0), Inches(0.3))
text(t, "the call chain", size=11.5, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
chain = [
    (f"encode()", fn(TOK, "encode"), "BOS, NFKC, split on whitespace"),
    (f"_encode_word()", fn(TOK, "_encode_word"), "runs by character class"),
    (f"_emit_alpha()", fn(TOK, "_emit_alpha"), "analyse, then look up both tokens"),
    (f"_emit_clitic()", fn(TOK, "_emit_clitic"), "kind='p' / 'e' picks the prefix range"),
]
x = 0.55
for name, ref, note in chain:
    box(s, x, 5.78, 3.02, 0.72,
        [(name, {"size": 11, "bold": True, "font": MONO, "color": INK}),
         (ref, {"size": 8, "font": MONO, "color": GREEN}),
         (note, {"size": 8.2, "color": MUTED})],
        fill=WHITE, line=GREEN, space=0)
    if x < 9.6:
        arrow(s, x + 3.06, 6.07, 0.13, 0.14, "right", GREEN)
    x += 3.24

footnote(s, "TokenizerOutput.tokens carries cleaned Arabic surface strings in parallel with the ids — "
            "the morphological metrics read those, never the vocabulary strings.")


# ==========================================================================
# 4 — clitics and pattern normalization
# ==========================================================================
s = blank(prs)
title(s, "Clitics — how the pattern is cleaned", "CAMeL bakes clitic surface characters into the pattern. They are stripped so a PAT token is a bare-stem template.",
      accent=INDIGO)

box(s, 0.55, 1.28, 12.22, 0.46,
    [("prc* and enc0 are FEATURE TAGS, not Arabic:  'wa_conj' → و  ·  'Al_det' → ال  ·  '3ms_dobj' → ه"
      "     translated by clitic_surface()  " + fn(BACK, "clitic_surface"),
      {"size": 10.5, "bold": True, "color": INDIGO})],
    fill=INDIGO_L, line=None, radius=0.1)

t = s.shapes.add_textbox(Inches(0.55), Inches(1.88), Inches(6.2), Inches(0.28))
text(t, f"normalize_pattern()  {fn(BACK,'normalize_pattern')}  — والكتاب, step by step",
     size=11, bold=True, color=INDIGO, align=PP_ALIGN.LEFT)

steps = [
    ("CAMeL raw pattern", "وَال1ُ2ّا3ِ", "prc2='wa_conj' → و,  prc0='Al_det' → ال", WHITE),
    ("strip prc3 (none)", "وَال1ُ2ّا3ِ", "outermost first", WHITE),
    ("strip prc2  و", "ال1ُ2ّا3ِ", "_strip_clitic_from_start skips the fatha too", INDIGO_L),
    ("strip prc1 (none)", "ال1ُ2ّا3ِ", "", WHITE),
    ("strip prc0  ال", "1ُ2ّا3ِ", "← the bare-stem template stored in vocab", GREEN_L),
]
y = 2.22
for label, pat, note, fill in steps:
    box(s, 0.55, y, 2.25, 0.42, label, fill=WHITE, line=SLATE_L, lw=0.7, size=9, align=PP_ALIGN.LEFT)
    box(s, 2.88, y, 1.95, 0.42, pat, fill=fill, line=INDIGO if fill is not WHITE else SLATE_L,
        lw=0.9, size=13, bold=True, color=INDIGO if fill is not WHITE else INK)
    box(s, 4.92, y, 3.55, 0.42, note, fill=WHITE, line=None, size=8.5, color=MUTED, align=PP_ALIGN.LEFT)
    y += 0.47

box(s, 0.55, 4.70, 7.92, 0.92,
    [("Order matters — outer to inner", {"size": 10.5, "bold": True, "color": ROSE}),
     ("CAMeL stacks prc3 (question) > prc2 (conjunction) > prc1 (preposition/future) > prc0 (article). "
      "Strip in that order or a later clitic partial-matches an earlier position.", {"size": 9, "color": MUTED})],
    fill=ROSE_L, line=ROSE, space=1)

box(s, 8.65, 1.88, 4.13, 3.74,
    [("the لِ + الـ contraction", {"size": 11, "bold": True, "color": ROSE}),
     ("Arabic writes one lam, not two:", {"size": 9, "color": MUTED}),
     ("لِ + الوَلَد  →  لِلوَلَد", {"size": 14, "bold": True}),
     ("", {"size": 3}),
     ("so the article's surface is ل, not ال, and a", {"size": 8.6, "color": MUTED}),
     ("literal strip of ال silently fails:", {"size": 8.6, "color": MUTED}),
     ("لِل1ِ2ا3ِ → ل1ِ2ا3ِ   ✗ stray lam", {"size": 10, "font": MONO, "color": ROSE}),
     ("لِل1ِ2ا3ِ → 1ِ2ا3ِ    ✓ handled", {"size": 10, "font": MONO, "color": GREEN}),
     ("", {"size": 3}),
     ("encode  strip_proclitics_from_start()", {"size": 8.6, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("decode  join_proclitics()", {"size": 8.6, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("the two MUST stay inverse", {"size": 9, "bold": True, "color": ROSE})],
    fill=WHITE, line=ROSE, space=1)

ex = [("الكتاب", "[CLITICP_ال]", "1ِ2ا3ِ"), ("لكتاب", "[CLITICP_ل]", "1ِ2ا3ِ"),
      ("للكتاب", "[CLITICP_ل] [CLITICP_ال]", "1ِ2ا3ِ"), ("بالكتاب", "[CLITICP_ب] [CLITICP_ال]", "1ِ2ا3ِ")]
t = s.shapes.add_textbox(Inches(0.55), Inches(5.72), Inches(7.9), Inches(0.28))
text(t, "four clitic shapes, one shared pattern token — that is the point of decoupling",
     size=10.5, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
x = 0.55
for w, cl, pat in ex:
    box(s, x, 6.06, 1.92, 0.82,
        [(w, {"size": 13, "bold": True}), (cl, {"size": 7.6, "color": INDIGO}),
         (f"[PAT_{pat}]", {"size": 8.4, "bold": True, "color": GREEN})],
        fill=WHITE, line=SLATE_L, space=0)
    x += 1.99


# ==========================================================================
# 5 — the # masked radical
# ==========================================================================
s = blank(prs)
title(s, "The masked radical  #  — a radical, not a missing field",
      "CAMeL writes # where a root consonant's surface form is not stable across the paradigm: the weak letters و / ي / ا and the hamza family.",
      accent=AMBER)

box(s, 0.55, 1.30, 6.05, 2.05,
    [("one lexeme, one root token", {"size": 11, "bold": True, "color": AMBER}),
     ("قال    root='ق.#.ل'   pattern='1ا3َ'", {"size": 10.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("يقول   root='ق.#.ل'   pattern='يَ1ُو3'", {"size": 10.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("قول    root='ق.#.ل'   pattern='1َوْ3ِ'", {"size": 10.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("أقوال  root='ق.#.ل'   pattern='أَ1ْوا3'", {"size": 10.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("", {"size": 3}),
     ("all four → [ROOT_ق#ل], the surface letter living in the pattern",
      {"size": 9, "color": MUTED})],
    fill=AMBER_L, line=AMBER, space=1)

box(s, 6.75, 1.30, 6.03, 2.05,
    [("# and the pattern are complementary", {"size": 11, "bold": True, "color": GREEN}),
     ("The masked radical's slot digit is ABSENT from the pattern;", {"size": 9, "color": MUTED, "align": PP_ALIGN.LEFT}),
     ("the letter that surfaces sits there as literal template material.", {"size": 9, "color": MUTED, "align": PP_ALIGN.LEFT}),
     ("", {"size": 4}),
     ("ق#ل + 1ا3َ  → قال        (slot 2 missing, ا literal)", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("ق#ل + يَ1ُو3 → يقول       (و literal)", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("", {"size": 3}),
     ("verified on 457 / 457 masked radicals — zero counter-examples",
      {"size": 9, "bold": True, "color": GREEN})],
    fill=GREEN_L, line=GREEN, space=1)

t = s.shapes.add_textbox(Inches(0.55), Inches(3.50), Inches(12.2), Inches(0.28))
text(t, f"why deletion was destructive — _dict_to_analysis()  {fn(BACK,'_dict_to_analysis')}",
     size=11, bold=True, color=ROSE, align=PP_ALIGN.LEFT)
box(s, 0.55, 3.84, 6.05, 1.30,
    [("deleting # ...", {"size": 10, "bold": True, "color": ROSE}),
     ("ق.#.ل → 'قل'  — 2 letters, below the 3-radical bar → rejected", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("and it RENUMBERS the survivors: ل moves from index 2 to 1,", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("so a pattern referring to slot 3 reads out of range", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("cost: 49.3 % of word occurrences → character path", {"size": 9.5, "bold": True, "color": ROSE})],
    fill=ROSE_L, line=ROSE, space=0)
box(s, 6.75, 3.84, 6.03, 1.30,
    [("... so radicals are counted structurally", {"size": 10, "bold": True, "color": GREEN}),
     ("radicals = [r for r in root.replace('_','.').split('.') if r]", {"size": 8.6, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("if len(radicals) < 3: return None", {"size": 8.6, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("WEAK_RADICAL_MARK is whitelisted by _is_arabic_root() and", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("guarded in naive_pattern_fill() so # never reaches output text", {"size": 9, "align": PP_ALIGN.LEFT})],
    fill=GREEN_L, line=GREEN, space=0)

box(s, 0.55, 5.32, 12.23, 1.35,
    [("# is lossy — know the trade-off", {"size": 11, "bold": True, "color": AMBER}),
     ("[ROOT_س#ر] covers both س-و-ر (سور, a wall) and س-ي-ر (سار, to march). "
      "[ROOT_ص#م] covers صوم and صام.", {"size": 10}),
     ("The (root, pattern) PAIR stays unambiguous — the pattern carries the realized letter — but the "
      "root token alone does not, which softens 'each ROOT is a semantic atom' for weak roots.", {"size": 9.5, "color": MUTED}),
     ("Canonicalizing # to a guessed radical was rejected: it merges exactly the same pairs while hiding "
      "the ambiguity behind a plausible-looking root.", {"size": 9.5, "color": MUTED})],
    fill=AMBER_L, line=AMBER, space=2)


# ==========================================================================
# 6 — path ②: the five gates that send a word to LIT
# ==========================================================================
s = blank(prs)
title(s, "Path ② — the five gates to the character fallback",
      "A word reaches [LIT_*] when any one of these fires. Four are in _dict_to_analysis(); the fifth is the vocabulary check in _emit_alpha().",
      accent=ROSE)

gates = [
    ("1", "no usable candidate", "CAMeL's backoff analyzer fires:  root='O', pattern='backoff'",
     "~33 % of word occurrences — no database entry", "_first_valid() returns None"),
    ("2", "fewer than 3 radicals", "في  root='ف.#'  ·  كم  root='ك.م'  ·  أي  root='#.#'",
     "genuinely bi-consonantal words and particles", "len(radicals) < 3"),
    ("3", "database marker", "تلفزيون · مريم  root='NTWS'   ·   دورهٔ  root='FOREIGN'",
     "loanwords, proper nouns, non-Arabic orthography", "root in ('NTWS','FOREIGN')"),
    ("4", "root is not Arabic", "هيكتور  root='Uٌٍ'  — an ASCII fragment from the database",
     "catch-all so junk never becomes a token", "_is_arabic_root()"),
    ("5", "root or pattern out of vocab", "القرآن  root='قر#' ✓ in vocab   pattern='1ُ2ْآنِ' ✗ OOV",
     "all-or-nothing: the in-vocab root is discarded too", "in self._vocab"),
]
y = 1.30
for num, name, example, note, guard in gates:
    box(s, 0.55, y, 0.42, 0.92, num, fill=ROSE, line=None, size=13, bold=True, color=WHITE, radius=0.3)
    box(s, 1.05, y, 2.55, 0.92, name, fill=ROSE_L, line=ROSE, lw=0.9, size=10.5, bold=True, color=ROSE)
    box(s, 3.68, y, 5.35, 0.92, example, fill=WHITE, line=SLATE_L, lw=0.7, size=10,
        align=PP_ALIGN.LEFT)
    box(s, 9.11, y, 3.67, 0.92,
        [(note, {"size": 8.8, "color": MUTED, "align": PP_ALIGN.LEFT}),
         (guard, {"size": 8.4, "font": MONO, "color": ROSE, "align": PP_ALIGN.LEFT})],
        fill=WHITE, line=SLATE_L, lw=0.7, space=1)
    y += 0.98

box(s, 0.55, 6.28, 12.23, 0.62,
    [("_emit_lit()  " + fn(TOK, "_emit_lit") +
      "   —   تلفزيون → [LIT_BEGIN] [CHAR_ت][CHAR_ل][CHAR_ف][CHAR_ز][CHAR_ي][CHAR_و][CHAR_ن] [LIT_END]",
      {"size": 11, "bold": True, "color": WHITE})],
    fill=ROSE, line=None, radius=0.1)


# ==========================================================================
# 7 — non-alpha runs and the run splitter
# ==========================================================================
s = blank(prs)
title(s, "Paths ③ ④ ⑤ — the word is not treated as atomic",
      "_encode_word() walks the whitespace word and cuts at every change of character class, so one word can emit several runs.",
      accent=AMBER)

box(s, 0.55, 1.30, 5.9, 2.35,
    [(f"_classify_char()   {fn(TOK,'_classify_char')}", {"size": 10.5, "bold": True, "font": MONO, "color": AMBER}),
     ("", {"size": 3}),
     ("alpha   35 ARABIC_LETTERS + 12 ARABIC_DIACRITICS", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("digit   0-9 and ٠-٩", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("punct   ASCII + ، ؛ ؟ « » … — – ـ", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("space   already removed by the outer split", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("other   everything else → <unk>", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT, "color": ROSE}),
     ("", {"size": 3}),
     ("ك→alpha   2→digit   !→punct   a→other   ى→other", {"size": 9, "color": MUTED})],
    fill=AMBER_L, line=AMBER, space=1)

box(s, 6.65, 1.30, 6.13, 2.35,
    [("runs, worked", {"size": 10.5, "bold": True, "color": AMBER}),
     ("\"يذهب\"   → [alpha 'يذهب']                     one run", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("\"2024م\"  → [digit '2024'] [alpha 'م']          two runs", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("\"عليكم!\" → [alpha 'عليكم'] [punct '!']         two runs", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("\"إلى\"    → [alpha 'إل'] [other 'ى']            two runs", {"size": 9.5, "font": MONO, "align": PP_ALIGN.LEFT, "color": ROSE}),
     ("", {"size": 4}),
     ("Each run is emitted independently — the analyser only ever sees", {"size": 9, "color": MUTED, "align": PP_ALIGN.LEFT}),
     ("the alpha runs, which is why the same chunking must be used by", {"size": 9, "color": MUTED, "align": PP_ALIGN.LEFT}),
     ("the corpus pre-pass (_extract_alpha_chunks) or the cache misses.", {"size": 9, "color": MUTED, "align": PP_ALIGN.LEFT})],
    fill=WHITE, line=AMBER, space=1)

box(s, 0.55, 3.85, 12.23, 1.35,
    [("known defect — ى (alef maqsura, U+0649) is not in ARABIC_LETTERS", {"size": 11.5, "bold": True, "color": ROSE}),
     ("It classifies as 'other', so it both SPLITS the word and becomes <unk>. "
      "normalize_arabic() folds the alef variants آ أ إ ٱ → ا but does not map ى → ي, so nothing upstream rescues it.",
      {"size": 9.5, "color": MUTED}),
     ("encode(\"إلى المستشفى\")  →  [LIT_BEGIN][CHAR_إ][CHAR_ل][LIT_END] <unk> [LIT_BEGIN][CHAR_ا]…[CHAR_ف][LIT_END] <unk>",
      {"size": 9, "font": MONO, "color": ROSE}),
     ("decode → 'إل ? المستشف ?'    ·    it also truncates the chunk before CAMeL ever sees it",
      {"size": 9.5, "bold": True, "color": ROSE})],
    fill=ROSE_L, line=ROSE, space=2)

box(s, 0.55, 5.42, 6.05, 1.25,
    [("③ digits glue back at decode", {"size": 10.5, "bold": True, "color": AMBER}),
     ("[DIGIT_2][DIGIT_0][DIGIT_2][DIGIT_4] would decode as four", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("separate words. decode() glues a digit onto the previous", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("output when it ends in a digit and no proclitic is buffered.", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("\"في عام 2024\" round-trips exactly.", {"size": 9, "bold": True, "color": GREEN, "align": PP_ALIGN.LEFT})],
    fill=WHITE, line=AMBER, space=0)
box(s, 6.75, 5.42, 6.03, 1.25,
    [("⑥ unknown clitic surface", {"size": 10.5, "bold": True, "color": ROSE}),
     ("_emit_clitic() looks up [CLITICP_x] or [CLITICE_x] by kind.", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("If that exact token is absent, the clitic is spelled out via", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("_emit_lit() instead of being dropped — the sixth path, and", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("the reason decode never loses a morpheme it cannot name.", {"size": 9, "align": PP_ALIGN.LEFT})],
    fill=WHITE, line=ROSE, space=0)


# ==========================================================================
# 8 — training: how the vocabulary is built
# ==========================================================================
s = blank(prs)
title(s, "train() — five steps, and what each produces",
      f"{fn(TOK,'train')}   ·   vocab_size is ignored on purpose: the size falls out of max_roots + max_patterns + fixed slots.",
      accent=INDIGO)

steps = [
    ("1", "_corpus_prepass()", fn(TOK, "_corpus_prepass"),
     "Analyse every distinct alpha chunk once, in batches of 256, and cache to disk.",
     "1,140,803 chunks → 506,101 analysed (44.4 %)"),
    ("2", "frequency tables", "",
     "Counters over accepted entries: root_freq, pat_freq, proclitic_freq, enclitic_freq.",
     "4,331 roots · 18,970 patterns · 10 proclitics · 13 enclitics"),
    ("3", "_build_vocab()", fn(TOK, "_build_vocab"),
     "Assemble the flat dict in fixed class order; roots and patterns freq-sorted, freq ≥ 2.",
     "8,265 tokens — 4,128 roots + 4,000 patterns + 137 fixed"),
    ("4", "_build_reconstruction()", fn(TOK, "_build_reconstruction"),
     "(root_id, pat_id) → inflected stem = diac minus clitic surfaces. Generator fills unseen pairs.",
     "108,149 entries"),
    ("5", "_build_metadata()", fn(TOK, "_build_metadata"),
     "Provenance: per-root and per-pattern id, freq, source and example words.",
     "vocab_metadata.json"),
]
y = 1.30
for num, name, ref, what, out in steps:
    box(s, 0.55, y, 0.42, 0.86, num, fill=INDIGO, line=None, size=13, bold=True, color=WHITE, radius=0.3)
    box(s, 1.05, y, 3.05, 0.86,
        [(name, {"size": 10.5, "bold": True, "font": MONO, "color": INDIGO}),
         (ref, {"size": 8, "font": MONO, "color": MUTED})],
        fill=INDIGO_L, line=INDIGO, lw=0.9, space=0)
    box(s, 4.18, y, 4.85, 0.86, what, fill=WHITE, line=SLATE_L, lw=0.7, size=9.2,
        align=PP_ALIGN.LEFT, color=INK)
    box(s, 9.11, y, 3.67, 0.86, out, fill=WHITE, line=INDIGO, lw=0.9, size=9.5, bold=True, color=INDIGO)
    y += 0.92

box(s, 0.55, 5.95, 6.05, 0.95,
    [("the cache is keyed on chunking", {"size": 10.5, "bold": True, "color": ROSE}),
     ("_extract_alpha_chunks() must match what _encode_word() does at", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("runtime, or every encode call misses and hits the analyser again.", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("Change one, change both.", {"size": 9, "bold": True, "color": ROSE, "align": PP_ALIGN.LEFT})],
    fill=ROSE_L, line=ROSE, space=0)
box(s, 6.75, 5.95, 6.03, 0.95,
    [("the budget binds on patterns, not roots", {"size": 10.5, "bold": True, "color": INDIGO}),
     ("Only 4,128 roots clear min_root_freq, so max_roots (10,000) never", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("binds. 15,219 patterns qualify and 4,000 are kept — admission goes", {"size": 9, "align": PP_ALIGN.LEFT}),
     ("26 % → 48 % as that budget rises from 500 to 4,000.", {"size": 9, "align": PP_ALIGN.LEFT})],
    fill=INDIGO_L, line=INDIGO, space=0)


# ==========================================================================
# 9 — vocabulary layout + the grid
# ==========================================================================
s = blank(prs)
title(s, "The vocabulary — contiguous ranges by class",
      "IDs are assigned in a fixed order so a class can be masked with a slice, not a set membership test.",
      accent=INDIGO)

segs = [("specials", "0–5", 6, SLATE, 1.00, "<pad> <s> </s> <unk> [LIT_BEGIN] [LIT_END]"),
        ("proclitics", "6–15", 10, TEAL, 1.35, "ال و ب ل ف ك س أ ما لا"),
        ("enclitics", "16–28", 13, TEAL, 1.35, "ه ها هم ي ك نا هما كم …"),
        ("chars", "29–75", 47, AMBER, 1.05, "Arabic letters + diacritics"),
        ("digits", "76–95", 20, AMBER, 0.95, "0-9 ٠-٩"),
        ("punct", "96–136", 41, AMBER, 1.00, "! . ، ؟ …"),
        ("ROOTS", "137–4264", 4128, INDIGO, 3.20, "freq-sorted, freq ≥ 2"),
        ("PATTERNS", "4265–8264", 4000, PURPLE, 2.05, "freq-sorted, freq ≥ 2")]
x = 0.55
for name, rng, n, col, w, note in segs:
    hot = name in ("ROOTS", "PATTERNS")
    box(s, x, 1.32, w, 1.05,
        [(name, {"size": 11 if hot else 9.5, "bold": True, "color": WHITE if hot else col}),
         (rng, {"size": 8.5, "color": RGBColor(0xE2, 0xE8, 0xF0) if hot else MUTED}),
         (f"{n:,}", {"size": 8.5, "color": RGBColor(0xE2, 0xE8, 0xF0) if hot else MUTED}),
         (note, {"size": 7.2, "color": RGBColor(0xE2, 0xE8, 0xF0) if hot else MUTED})],
        fill=col if hot else WHITE, line=col, lw=1.6 if hot else 1.0, space=0)
    x += w + 0.07

box(s, 0.55, 2.55, 12.22, 0.44,
    [("emission is all-or-nothing per word: [ROOT_x] and [PAT_y] must BOTH be present, "
      "or the word — root included — goes to the character path",
      {"size": 11, "bold": True, "color": INDIGO})],
    fill=INDIGO_L, line=None, radius=0.12)

t = s.shapes.add_textbox(Inches(0.55), Inches(3.10), Inches(12.2), Inches(0.3))
text(t, "why two tokens beat one entry — the (root × pattern) grid", size=12, bold=True,
     color=INK, align=PP_ALIGN.LEFT)

box(s, 0.55, 3.48, 6.05, 2.55,
    [("[ROOT_ذهب]  ×  58 patterns", {"size": 11, "bold": True, "color": INDIGO}),
     ("[PAT_1َ2َ3]     → ذهب", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[PAT_1ا2ِ3]     → ذاهب", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[PAT_يَ1ْ2َ3]    → يذهب", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[PAT_تَ1ْ2َ3]    → تذهب", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[PAT_مَ1ْ2َ3]    → مذهب", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[PAT_مَ1ا2ِ3]   → مذاهب", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("…52 more", {"size": 9, "color": MUTED})],
    fill=INDIGO_L, line=INDIGO, space=0)
box(s, 6.75, 3.48, 6.03, 2.55,
    [("[PAT_يَ1ْ2َ3]  ×  360 roots", {"size": 11, "bold": True, "color": PURPLE}),
     ("[ROOT_جمع] → يجمع        [ROOT_علم] → يعلم", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[ROOT_عمل] → يعمل        [ROOT_قطع] → يقطع", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[ROOT_دفع] → يدفع        [ROOT_حفظ] → يحفظ", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("[ROOT_ظهر] → يظهر        [ROOT_شهد] → يشهد", {"size": 10, "font": MONO, "align": PP_ALIGN.LEFT}),
     ("…352 more", {"size": 9, "color": MUTED}),
     ("", {"size": 4}),
     ("one shared token doing the grammatical work for every verb",
      {"size": 9, "color": MUTED})],
    fill=PURPLE_L, line=PURPLE, space=0)

box(s, 0.55, 6.20, 12.23, 0.62,
    [("58 + 360 = 418 surface forms addressable with two tokens instead of 418 separate vocabulary entries",
      {"size": 12.5, "bold": True, "color": WHITE})],
    fill=INDIGO, line=None, radius=0.12)


# ==========================================================================
# 10 — decode: the state machine
# ==========================================================================
s = blank(prs)
title(s, "decode() — one left-to-right pass, four pieces of state",
      f"{fn(TOK,'decode')}   ·   dispatch is on the token's string PREFIX, which is why proclitics and enclitics have distinct ranges.",
      accent=GREEN)

st = [("out", "finalized words, in order"), ("clitic_prefix", "proclitics buffered for the next word"),
      ("pending_root", "a ROOT that has not met its PAT yet"), ("lit_buffer", "chars between LIT_BEGIN/END")]
x = 0.55
for n, d in st:
    box(s, x, 1.30, 3.02, 0.68, [(n, {"size": 11, "bold": True, "font": MONO, "color": GREEN}),
                                 (d, {"size": 8.5, "color": MUTED})],
        fill=GREEN_L, line=GREEN, space=0)
    x += 3.09

hdr = [("token prefix", 3.10), ("branch", 5.55), ("state after", 3.57)]
x = 0.55
for h_, w_ in hdr:
    box(s, x, 2.14, w_, 0.38, h_, fill=SLATE, line=None, size=10, bold=True, color=WHITE,
        radius=0.04, align=PP_ALIGN.LEFT)
    x += w_
rows = [
    ("[CLITICP_*]", "buffer it — belongs to the word that follows", "clitic_prefix.append(...)"),
    ("[CLITICE_*]", "attach to the word just emitted", "out[-1] += surface"),
    ("[ROOT_*]", "dump any orphan root, then hold this one", "pending_root = 'ذهب'"),
    ("[PAT_*]", "_reconstruct(root_id, pat_id) → flush_word()", "out += join_proclitics(...) + stem"),
    ("[LIT_BEGIN] / [CHAR_*] / [LIT_END]", "accumulate literal characters, then flush", "lit_buffer"),
    ("[DIGIT_*]", "glue onto the previous output if it ends in a digit", "2 0 2 4 → 2024"),
    ("[PUNCT_*] / <unk>", "flush as its own word / emit '?'", "—"),
]
y = 2.56
for i, (a_, b_, c_) in enumerate(rows):
    fill = WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF5, 0xF9)
    hot = a_ == "[PAT_*]"
    x = 0.55
    for val, w_ in zip((a_, b_, c_), (3.10, 5.55, 3.57)):
        box(s, x, y, w_, 0.44, val, fill=GREEN_L if hot else fill, line=SLATE_L, lw=0.7,
            size=9.5, bold=hot, color=GREEN if hot else INK, align=PP_ALIGN.LEFT, radius=0.02)
        x += w_
    y += 0.44

box(s, 0.55, 5.75, 6.05, 1.10,
    [("edge cases the machine absorbs", {"size": 10.5, "bold": True, "color": GREEN}),
     ("PAT with no pending ROOT → naive_pattern_fill('', pat) returns '' → the word vanishes silently",
      {"size": 8.8, "align": PP_ALIGN.LEFT}),
     ("ROOT with no following PAT → dump_orphan_root() emits the bare root letters",
      {"size": 8.8, "align": PP_ALIGN.LEFT}),
     ("trailing buffered proclitics at EOS → flushed as a standalone word",
      {"size": 8.8, "align": PP_ALIGN.LEFT})],
    fill=WHITE, line=GREEN, space=0)
box(s, 6.75, 5.75, 6.03, 1.10,
    [("why prc and enc need distinct prefixes", {"size": 10.5, "bold": True, "color": ROSE}),
     ("With one [CLITIC_*] range the decoder cannot tell whether a ك attaches",
      {"size": 8.8, "align": PP_ALIGN.LEFT}),
     ("leftward or rightward — ka_prep 'like' vs 2ms_pron 'your' are the same",
      {"size": 8.8, "align": PP_ALIGN.LEFT}),
     ("surface. Distinct ranges make it a token-type decision, not a guess.",
      {"size": 8.8, "align": PP_ALIGN.LEFT})],
    fill=ROSE_L, line=ROSE, space=0)


# ==========================================================================
# 11 — the three-tier resolver
# ==========================================================================
s = blank(prs)
title(s, "_reconstruct() — three tiers, tried in order",
      f"{fn(TOK,'_reconstruct')}   ·   turning (root_id, pat_id) back into an Arabic surface form.",
      accent=GREEN)

tiers = [
    ("TIER 1", "lookup table", "self._reconstruction[(root_id, pat_id)]",
     "108,149 entries built at training time from the corpus, storing the most frequent realization.",
     "hits ~99 % of LLM emissions — the model was trained on exactly this distribution", "O(1)", GREEN, GREEN_L),
    ("TIER 2", "CAMeL Generator", "backend.generate(root, pattern)",
     "The Generator takes a lemma + features, which we do not have. So the server naive-fills, re-analyses "
     "that rough form, and returns the stem of whichever analysis matches the target (root, bare pattern).",
     "exploits the analyser as a lossy inverse of the generator — correct on weak roots", "~1 ms, cached", AMBER, AMBER_L),
    ("TIER 3", "naive substitution", "naive_pattern_fill(root, pattern)",
     "Substitute root letters into the pattern's slot digits; everything else is carried verbatim.",
     "wrong on weak roots — قَوَلَ instead of قال — so it is the safety net, never the plan", "O(L)", ROSE, ROSE_L),
]
y = 1.30
for name, what, call, how, why, cost, col, light in tiers:
    box(s, 0.55, y, 1.15, 1.62, name, fill=col, line=None, size=11, bold=True, color=WHITE, radius=0.18)
    box(s, 1.78, y, 11.0, 1.62,
        [(f"{what}      {call}", {"size": 11, "bold": True, "color": col, "align": PP_ALIGN.LEFT}),
         (how, {"size": 9.5, "align": PP_ALIGN.LEFT}),
         (why, {"size": 9.5, "color": MUTED, "align": PP_ALIGN.LEFT}),
         (f"cost: {cost}", {"size": 9, "bold": True, "color": col, "align": PP_ALIGN.LEFT})],
        fill=light, line=col, space=2)
    y += 1.72
    if y < 6.0:
        arrow(s, 6.6, y - 0.09, 0.16, 0.08, "down", MUTED)

box(s, 0.55, 6.50, 12.23, 0.62,
    [("reconstruction[(536, 4284)] = 'يذهب'      ·      reconstruction[(176, 4436)] = 'قال'      ·      "
      "reconstruction[(197, 4377)] = 'ولد'",
      {"size": 12, "bold": True, "color": WHITE, "font": MONO})],
    fill=SLATE, line=None, radius=0.1)


# ==========================================================================
# 12 — function map
# ==========================================================================
s = blank(prs)
title(s, "Function map — where to look when reading the code",
      "Line numbers are resolved when this deck is generated, so regenerating re-syncs them.",
      accent=PURPLE)

groups = [
    ("ENCODE", TEAL, TEAL_L, [
        ("encode()", TOK, "encode", "BOS/EOS, NFKC, whitespace split, padding & truncation"),
        ("_encode_word()", TOK, "_encode_word", "cut into runs by character class"),
        ("_classify_char()", TOK, "_classify_char", "alpha / digit / punct / space / other"),
        ("_emit_alpha()", TOK, "_emit_alpha", "analyse → vocab check → ROOT+PAT, else LIT"),
        ("_emit_clitic()", TOK, "_emit_clitic", "kind 'p'/'e' selects the prefix range"),
        ("_emit_lit()", TOK, "_emit_lit", "the character fallback"),
    ]),
    ("ANALYSE", AMBER, AMBER_L, [
        ("MorphAnalyzer.analyze()", BACK, "analyze", "LRU-cached; analyze_many() batches 256 per round-trip"),
        ("_first_valid()", BACK, "_first_valid", "walk candidates in MLE order, return the first accepted"),
        ("_dict_to_analysis()", BACK, "_dict_to_analysis", "the four gates + normalization → Analysis"),
        ("normalize_pattern()", BACK, "normalize_pattern", "strip clitic surfaces → bare-stem template"),
        ("strip_proclitics_from_start()", BACK, "strip_proclitics_from_start", "outer-to-inner, handles لِ+الـ"),
        ("clitic_surface()", BACK, "clitic_surface", "CAMeL feature tag → Arabic surface"),
    ]),
    ("TRAIN", INDIGO, INDIGO_L, [
        ("train()", TOK, "train", "five steps; vocab_size is ignored"),
        ("_corpus_prepass()", TOK, "_corpus_prepass", "batched analysis + on-disk cache"),
        ("_build_vocab()", TOK, "_build_vocab", "deterministic ID order by class"),
        ("_build_reconstruction()", TOK, "_build_reconstruction", "(root_id, pat_id) → inflected stem"),
        ("_strip_clitic_surfaces()", TOK, "_strip_clitic_surfaces", "diac minus clitics = inflected stem"),
    ]),
    ("DECODE", GREEN, GREEN_L, [
        ("decode()", TOK, "decode", "the state machine"),
        ("_reconstruct()", TOK, "_reconstruct", "lookup → generator → naive"),
        ("join_proclitics()", TOK, "join_proclitics", "inverse of strip_proclitics_from_start"),
        ("naive_pattern_fill()", BACK, "naive_pattern_fill", "tier 3; guards against emitting #"),
    ]),
]
x = 0.55
for name, col, light, items in groups:
    box(s, x, 1.28, 3.02, 0.34, name, fill=col, line=None, size=10, bold=True, color=WHITE, radius=0.12)
    y = 1.68
    for label, path, sym, note in items:
        box(s, x, y, 3.02, 0.80,
            [(label, {"size": 9.2, "bold": True, "font": MONO, "color": col, "align": PP_ALIGN.LEFT}),
             (fn(path, sym), {"size": 7.6, "font": MONO, "color": MUTED, "align": PP_ALIGN.LEFT}),
             (note, {"size": 8, "color": INK, "align": PP_ALIGN.LEFT})],
            fill=light, line=col, lw=0.8, space=0)
        y += 0.86
    x += 3.09

footnote(s, "The bridge (araroopat_bridge.py) and the server (tools/araroopat_camel_server.py) sit between "
            "ANALYSE and CAMeL — NDJSON over stdin/stdout, one outstanding request, fail-loud on EOF or timeout.")


for p in audit(prs):
    print("  !", p)

prs.save(OUT)
print("wrote", OUT, OUT.stat().st_size, "bytes,", len(prs.slides._sldIdLst), "slides")
